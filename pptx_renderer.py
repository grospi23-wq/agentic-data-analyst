"""
pptx_renderer.py
----------------
Bulletproof Executive PPTX Renderer for Agentic Data Analyst.
Final Version: Integrated Auto-Scaling and Overflow Protection.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
import logfire
from pathlib import Path
from typing import Optional

from schema import PresentationSpec, SlideContent

# --- High-End Executive Palette ---
COLORS = {
    "BG": RGBColor(0x0A, 0x19, 0x2F),         
    "ACCENT": RGBColor(0x64, 0xFF, 0xDA),     
    "TEXT_MAIN": RGBColor(0xCC, 0xD6, 0xF6),  
    "TEXT_BOLD": RGBColor(0xFF, 0xFF, 0xFF),  
    "BOX_BG": RGBColor(0x11, 0x22, 0x40),     
    "GREEN": RGBColor(0x10, 0xB9, 0x81),      
    "RED": RGBColor(0xF4, 0x3F, 0x5E),        
    "GRAY": RGBColor(0x4B, 0x55, 0x63)        
}

W = Inches(13.33)
H = Inches(7.5)

def _set_dark_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["BG"]

def _apply_text_style(tf, size: int, color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT, auto_size=False):
    """Systematic font application with overflow protection."""
    if auto_size:
        tf.word_wrap = True
    
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color

def _add_executive_footer(slide, page_num: int):
    tag_box = slide.shapes.add_textbox(Inches(0.5), H - Inches(0.5), Inches(4), Inches(0.3))
    tag_box.text_frame.text = "STRATEGIC ANALYSIS | CONFIDENTIAL"
    _apply_text_style(tag_box.text_frame, 10, COLORS["GRAY"])
    
    num_box = slide.shapes.add_textbox(W - Inches(1.5), H - Inches(0.5), Inches(1), Inches(0.3))
    num_box.text_frame.text = f"PAGE {page_num}"
    _apply_text_style(num_box.text_frame, 10, COLORS["GRAY"], align=PP_ALIGN.RIGHT)

def _add_accent_header(slide, title_text: str, color: Optional[RGBColor] = None):
    accent_color = color if color else COLORS["ACCENT"]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    # Title Protection Logic: Dynamic font size based on length
    font_size = 34
    if len(title_text) > 40: font_size = 28
    if len(title_text) > 60: font_size = 24

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), W - Inches(1.5), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title_text
    _apply_text_style(tf, font_size, COLORS["TEXT_BOLD"], bold=True)

def _add_native_chart(slide, chart_spec, x, y, cx, cy):
    """Injects a native, editable PowerPoint chart customized for the dark theme."""
    try:
        chart_data = CategoryChartData()
        points = [dp for dp in chart_spec.data_points[:7] if dp.value is not None]
        if not points:
            logfire.warn("Chart skipped — no valid data points for '{t}'", t=chart_spec.title)
            return None
        chart_data.categories = [dp.label for dp in points]
        chart_data.add_series(chart_spec.title, [dp.value for dp in points])

        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE
        }
        xl_chart_type = chart_type_map.get(chart_spec.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        graphic_frame = slide.shapes.add_chart(
            xl_chart_type, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart

        # Dark theme styling for the chart text
        chart.font.color.rgb = COLORS["TEXT_MAIN"]
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.color.rgb = COLORS["TEXT_MAIN"]

        return chart
    except Exception as e:
        logfire.error("Chart injection failed for '{t}': {err}", t=chart_spec.title, err=str(e))
        return None

def _populate_title_slide(slide, spec: PresentationSpec):
    
    _set_dark_bg(slide)
    
    # Cinematic line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, W/2 - Inches(1), H/2 - Inches(1.5), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["ACCENT"]
    line.line.fill.background()
    
    # Title Scaling for long titles
    main_title = spec.deck_title.upper()
    title_font_size = 52
    if len(main_title) > 30: title_font_size = 42
    if len(main_title) > 50: title_font_size = 32

    title_box = slide.shapes.add_textbox(Inches(0.5), H/2 - Inches(1), W - Inches(1), Inches(2))
    tf = title_box.text_frame
    tf.text = main_title
    _apply_text_style(tf, title_font_size, COLORS["ACCENT"], bold=True, align=PP_ALIGN.CENTER, auto_size=True)
    
    sub_box = slide.shapes.add_textbox(Inches(1), H/2 + Inches(1.2), W - Inches(2), Inches(0.8))
    tf_sub = sub_box.text_frame
    tf_sub.text = f"{spec.subtitle} | EXECUTIVE SUMMARY"
    _apply_text_style(tf_sub, 20, COLORS["TEXT_MAIN"], align=PP_ALIGN.CENTER, auto_size=True)

def _render_finding_slide(slide, slide_spec: SlideContent, page_num: int):

    _add_executive_footer(slide, page_num)
    
    # Check if a chart was generated
    has_chart = getattr(slide_spec, "chart", None) is not None
    
    if slide_spec.metric_callout:
        # Dynamic placement: shorter and higher if a chart is sharing the right side
        box_y = Inches(1.2) if has_chart else Inches(1.5)
        box_h = Inches(1.5) if has_chart else Inches(2.2)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, W - Inches(4.2), box_y, Inches(3.8), box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS["BOX_BG"]
        box.line.color.rgb = COLORS["ACCENT"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = slide_spec.metric_callout
        # Callout Font Scaling
        c_font_size = 32
        if len(slide_spec.metric_callout) > 30: c_font_size = 24
        if len(slide_spec.metric_callout) > 50: c_font_size = 18
        
        _apply_text_style(tf, c_font_size, COLORS["ACCENT"], bold=True, align=PP_ALIGN.CENTER)

    if has_chart:
        # Dynamic placement: below the callout, or taking the full right space
        chart_y = Inches(2.9) if slide_spec.metric_callout else Inches(1.5)
        chart_h = Inches(3.8) if slide_spec.metric_callout else Inches(5.0)
        _add_native_chart(
            slide, 
            slide_spec.chart, 
            x=W - Inches(4.5), 
            y=chart_y, 
            cx=Inches(4.2), 
            cy=chart_h
        )

    # Core Body with wrap protection (Left side)
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.0), Inches(5.2))
    tf = body_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(slide_spec.body_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.space_before = Pt(14)
        _apply_text_style(tf, 18, COLORS["TEXT_MAIN"])

def _render_action_plan(slide, slide_spec: SlideContent, page_num: int):
    _set_dark_bg(slide)
    _add_accent_header(slide, "STRATEGIC EXECUTION ROADMAP")
    _add_executive_footer(slide, page_num)
    
    for i, action in enumerate(slide_spec.body_bullets[:4]):
        top_offset = Inches(1.6 + (i * 1.35))
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top_offset + Inches(0.1), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS["ACCENT"]
        circle.line.fill.background()
        
        c_tf = circle.text_frame
        c_tf.text = str(i + 1)
        _apply_text_style(c_tf, 18, COLORS["BG"], bold=True, align=PP_ALIGN.CENTER)
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), top_offset, Inches(11.4), Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["BOX_BG"]
        bar.line.color.rgb = COLORS["GRAY"]
        
        b_tf = bar.text_frame
        b_tf.word_wrap = True
        b_tf.text = action
        # Scaling for long action items
        a_font_size = 18
        if len(action) > 100: a_font_size = 14
        _apply_text_style(b_tf, a_font_size, COLORS["TEXT_BOLD"], bold=True)

def render_pptx(spec: PresentationSpec, output_path: str) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    for i, slide_spec in enumerate(spec.slides):
        page_idx = i + 1
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            if slide_spec.slide_type == "title":
                _populate_title_slide(slide, spec)
            elif slide_spec.slide_type == "action_plan":
                _render_action_plan(slide, slide_spec, page_idx)
            else:
                # For Findings, Hypothesis, Limitations, etc.
                header_color = None
                if slide_spec.slide_type == "hypothesis":
                    v = spec.hypothesis_verdict.upper()
                    if "REJECTED" in v: header_color = COLORS["RED"]
                    elif "CONFIRMED" in v: header_color = COLORS["GREEN"]

                _set_dark_bg(slide)
                _add_accent_header(slide, slide_spec.title, color=header_color)
                _render_finding_slide(slide, slide_spec, page_idx)

            # Handle speaker notes
            if slide_spec.speaker_notes:
                try:
                    notes_slide = prs.slides[-1].notes_slide
                    if notes_slide and notes_slide.notes_text_frame:
                        notes_slide.notes_text_frame.text = slide_spec.speaker_notes
                except Exception as e:
                    logfire.warn("Speaker notes failed for slide {n}: {err}", n=page_idx, err=str(e))

        except Exception as e:
            logfire.error("Slide {n} render failed ({t}): {err}", n=page_idx, t=slide_spec.slide_type, err=str(e))
            print(f"⚠️  Slide {page_idx} ({slide_spec.slide_type}) skipped due to render error: {e}")

    final_path = Path(output_path)
    prs.save(str(final_path))
    return final_path